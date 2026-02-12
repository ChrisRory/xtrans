

import sys
from pathlib import Path
from PIL import Image, ImageDraw
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import tempfile
import io
import tkinter as tk
from tkinter import filedialog
import random


def remove_watermark_from_pdf(input_pdf_path: str, output_pptx_path: str, dpi: int = 100):
    """
    从PDF中移除右下角水印并输出为PPTX

    参数:
        input_pdf_path: 输入PDF文件路径
        output_pptx_path: 输出PPTX文件路径
        dpi: 转换分辨率 (默认100)
    """

    print(f"正在处理: {input_pdf_path}")
    print(f"目标DPI: {dpi}")

    # 创建临时目录存储处理后的图片
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)

        # 步骤1: 将PDF转换为图片 (100 DPI)
        print(f"步骤1: 转换PDF为{dpi} DPI图片...")
        images = convert_from_path(input_pdf_path, dpi=dpi)
        print(f"共 {len(images)} 页")

        processed_images = []

        # 步骤2: 处理每一页
        for page_num, image in enumerate(images, 1):
            print(f"步骤2: 处理第 {page_num}/{len(images)} 页...")

            # 获取图片尺寸
            width, height = image.size
            print(f"  页面尺寸: {width} x {height} 像素")

            # 定义水印区域 (右下角 150x35)
            watermark_width = 150
            watermark_height = 35

            # 计算矩形坐标 (左上角和右下角)
            x1 = width - watermark_width
            y1 = height - watermark_height
            x2 = width
            y2 = height

            print(f"  覆盖区域: ({x1}, {y1}) 到 ({x2}, {y2})")

            # 获取3个参考像素的颜色
            pixel_bottom_right = image.getpixel((width - 1, height - 1))
            pixel_top = image.getpixel((width - 1, y1 - 1)) if y1 > 0 else pixel_bottom_right
            pixel_left = image.getpixel((x1 - 1, height - 1)) if x1 > 0 else pixel_bottom_right

            print(f"  参考颜色 - 右下: {pixel_bottom_right}, 上: {pixel_top}, 左: {pixel_left}")

            # 创建绘图对象
            draw = ImageDraw.Draw(image)

            # 使用随机混合颜色填充每个像素
            for x in range(x1, x2):
                for y in range(y1, y2):
                    # 随机选择使用哪个参考颜色或混合颜色
                    choice = random.randint(0, 2)
                    if choice == 0:
                        color = pixel_bottom_right
                    elif choice == 1:
                        color = pixel_top
                    else:
                        color = pixel_left

                    # 可选：添加轻微的颜色变化以使其更自然
                    if isinstance(color, tuple) and len(color) >= 3:
                        r = max(0, min(255, color[0] + random.randint(-2, 2)))
                        g = max(0, min(255, color[1] + random.randint(-2, 2)))
                        b = max(0, min(255, color[2] + random.randint(-2, 2)))
                        color = (r, g, b) if len(color) == 3 else (r, g, b, color[3])

                    draw.point((x, y), fill=color)

            # 转换为RGB模式
            if image.mode != 'RGB':
                image = image.convert('RGB')

            processed_images.append(image)
            print(f"  已处理页面 {page_num}")

        # 步骤3: 将处理后的图片转换为PPTX
        print(f"步骤3: 创建PPTX演示文稿...")
        prs = Presentation()

        # 设置幻灯片尺寸为标准16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)  # 10 / 16 * 9 = 5.625

        for page_num, image in enumerate(processed_images, 1):
            print(f"  添加幻灯片 {page_num}/{len(processed_images)}...")

            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]  # 6是空白布局
            slide = prs.slides.add_slide(blank_slide_layout)

            # 将图片保存到内存中
            img_stream = io.BytesIO()
            image.save(img_stream, format='PNG')
            img_stream.seek(0)

            # 图片填满整个幻灯片（因为PDF内容已经是16:9）
            slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # 保存PPTX
        prs.save(output_pptx_path)
        print(f"✓ 完成! 输出文件: {output_pptx_path}")

        # 显示文件大小
        input_size = Path(input_pdf_path).stat().st_size / 1024 / 1024
        output_size = Path(output_pptx_path).stat().st_size / 1024 / 1024
        print(f"原始文件: {input_size:.2f} MB")
        print(f"输出文件: {output_size:.2f} MB")


def select_pdf_file():
    """使用GUI选择PDF文件"""
    root = tk.Tk()
    root.withdraw()  # 隐藏主窗口

    file_path = filedialog.askopenfilename(
        title="选择PDF文件",
        filetypes=[("PDF files", "*.pdf"), ("All files", "*.*")]
    )

    return file_path


def main():
    """主程序入口"""
    # 如果有命令行参数，使用命令行模式
    if len(sys.argv) >= 2:
        input_pdf = sys.argv[1]

        # 如果没有指定输出文件名，自动生成
        if len(sys.argv) >= 3:
            output_pptx = sys.argv[2]
        else:
            input_path = Path(input_pdf)
            output_pptx = str(input_path.parent / f"{input_path.stem}.pptx")

        # DPI参数 (默认100)
        dpi = int(sys.argv[3]) if len(sys.argv) >= 4 else 100
    else:
        # 否则使用GUI模式
        print("请选择PDF文件...")
        input_pdf = select_pdf_file()

        if not input_pdf:
            print("未选择文件，退出。")
            sys.exit(0)

        # 自动生成输出文件名（相同目录，相同文件名，.pptx扩展名）
        input_path = Path(input_pdf)
        output_pptx = str(input_path.parent / f"{input_path.stem}.pptx")

        # 默认DPI
        dpi = 100

    # 检查输入文件是否存在
    if not Path(input_pdf).exists():
        print(f"错误: 文件不存在 - {input_pdf}")
        sys.exit(1)

    try:
        remove_watermark_from_pdf(input_pdf, output_pptx, dpi)
    except Exception as e:
        print(f"错误: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
