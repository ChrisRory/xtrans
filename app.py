import gradio as gr
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw
import random
import tempfile
import os
import io

def remove_watermark_and_convert(pdf_file, dpi=100, progress=gr.Progress()):
    """
    Gradio handler to convert PDF to PPTX with watermark removal.
    """
    if pdf_file is None:
        return None

    # --- Step 1: Convert PDF to Images ---
    progress(0, desc="Starting conversion...")
    
    try:
        # Gradio passes a temp file object. We use .name to get the path.
        print(f"Processing: {pdf_file.name}")
        images = convert_from_path(pdf_file.name, dpi=dpi)
    except Exception as e:
        raise gr.Error(f"Error reading PDF. Is Poppler installed? Details: {e}")

    total_pages = len(images)
    processed_images = []

    # --- Step 2: Remove Watermark ---
    for i, image in enumerate(images):
        # Update progress bar (0% to 50%)
        progress((i / total_pages) * 0.5, desc=f"Removing watermark: Page {i+1}/{total_pages}")

        width, height = image.size
        
        # Define watermark area (Bottom Right 150x35)
        watermark_width = 150
        watermark_height = 35
        
        x1 = width - watermark_width
        y1 = height - watermark_height
        x2 = width
        y2 = height

        # Get reference pixels
        pixel_bottom_right = image.getpixel((width - 1, height - 1))
        pixel_top = image.getpixel((width - 1, y1 - 1)) if y1 > 0 else pixel_bottom_right
        pixel_left = image.getpixel((x1 - 1, height - 1)) if x1 > 0 else pixel_bottom_right

        draw = ImageDraw.Draw(image)

        # Pixel manipulation loop (Your original logic)
        for x in range(x1, x2):
            for y in range(y1, y2):
                choice = random.randint(0, 2)
                if choice == 0:
                    color = pixel_bottom_right
                elif choice == 1:
                    color = pixel_top
                else:
                    color = pixel_left

                # Add slight noise
                if isinstance(color, tuple) and len(color) >= 3:
                    r = max(0, min(255, color[0] + random.randint(-2, 2)))
                    g = max(0, min(255, color[1] + random.randint(-2, 2)))
                    b = max(0, min(255, color[2] + random.randint(-2, 2)))
                    color = (r, g, b) if len(color) == 3 else (r, g, b, color[3])

                draw.point((x, y), fill=color)

        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        processed_images.append(image)

    # --- Step 3: Create PPTX ---
    progress(0.5, desc="Generating PowerPoint slides...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for i, image in enumerate(processed_images):
        # Update progress bar (50% to 100%)
        progress(0.5 + ((i / total_pages) * 0.5), desc=f"Adding slide {i+1}/{total_pages}")
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Save image to memory buffer for PPTX insertion
        img_stream = io.BytesIO()
        image.save(img_stream, format='PNG')
        img_stream.seek(0)

        slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # --- Step 4: Save and Return ---
    # Create a temporary file path for the output
    # We use the original filename stem to name the output
    original_stem = os.path.splitext(os.path.basename(pdf_file.name))[0]
    # Handle weird Gradio temp names if necessary, but usually safe to just make a new temp
    output_dir = tempfile.gettempdir()
    output_path = os.path.join(output_dir, f"{original_stem}_converted.pptx")
    
    prs.save(output_path)
    
    progress(1.0, desc="Done!")
    return output_path

# --- Build the Gradio UI ---
with gr.Blocks(title="PDF to PPTX Cleaner") as demo:
    gr.Markdown("# 📄 PDF to PPTX Watermark Remover")
    gr.Markdown("Upload a PDF to remove the bottom-right watermark and convert it to PowerPoint.")
    
    with gr.Row():
        with gr.Column():
            # Input Section
            input_file = gr.File(label="Upload PDF", file_types=[".pdf"], type="filepath")
            
            with gr.Accordion("Advanced Settings", open=True):
                dpi_slider = gr.Slider(
                    minimum=72, 
                    maximum=200, 
                    value=100, 
                    step=10, 
                    label="Quality (DPI)", 
                    info="Higher DPI = Better quality but slower processing."
                )
            
            convert_btn = gr.Button("🚀 Start Conversion", variant="primary")

        with gr.Column():
            # Output Section
            output_file = gr.File(label="Download PowerPoint")

    # Event Listener
    convert_btn.click(
        fn=remove_watermark_and_convert,
        inputs=[input_file, dpi_slider],
        outputs=output_file
    )

if __name__ == "__main__":
    # Support for Render/Docker Port mapping
    port = int(os.environ.get("PORT", 7860))
    demo.launch(server_name="0.0.0.0", server_port=port)
